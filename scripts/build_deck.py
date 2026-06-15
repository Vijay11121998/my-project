#!/usr/bin/env python3
"""Generate the client presentation (.pptx) for the Bank Reconciliation Agent.

Mirrors docs/reconciliation-agent-deck.md. Run: python3 scripts/build_deck.py
Output: docs/reconciliation-agent-deck.pptx
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- Theme ----------------------------------------------------------------
NAVY = RGBColor(0x0E, 0x2A, 0x47)      # deep navy - trust
TEAL = RGBColor(0x12, 0x8C, 0x7E)      # healthcare teal - accent
LIGHT = RGBColor(0xF4, 0xF7, 0xF9)     # light background
INK = RGBColor(0x1B, 0x2A, 0x33)       # body text
MUTE = RGBColor(0x5A, 0x6B, 0x76)      # muted text
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_BG = RGBColor(0xE7, 0xF1, 0xEF)  # pale teal panel

FONT = "Calibri"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)  # 16:9

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


# ---- helpers --------------------------------------------------------------
def add_slide():
    return prs.slides.add_slide(BLANK)


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def rect(slide, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    fill(sh, color)
    sh.shadow.inherit = False
    return sh


def textbox(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def set_run(run, text, size, color, bold=False, italic=False, font=FONT):
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def header(slide, title, kicker="MONTH-END RECONCILIATION AGENT"):
    """Standard content-slide header: teal bar + title."""
    rect(slide, 0, 0, EMU_W, Inches(0.18), TEAL)
    tb, tf = textbox(slide, Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.35))
    p = tf.paragraphs[0]
    set_run(p.add_run(), kicker, 11, TEAL, bold=True)
    tb2, tf2 = textbox(slide, Inches(0.6), Inches(0.7), Inches(12.1), Inches(0.9))
    p2 = tf2.paragraphs[0]
    set_run(p2.add_run(), title, 30, NAVY, bold=True)


def bullets(slide, items, top=Inches(1.9), left=Inches(0.7),
            width=Inches(12.0), size=18, gap=10):
    """items: list of (text, level) or str."""
    tb, tf = textbox(slide, left, top, width, Inches(5.0))
    first = True
    for it in items:
        text, level = (it, 0) if isinstance(it, str) else it
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(gap)
        marker = "•  " if level == 0 else "–  "
        set_run(p.add_run(), marker, size, TEAL, bold=True)
        set_run(p.add_run(), text, size - (2 if level else 0), INK)


# ---- Slide 1: Title -------------------------------------------------------
s = add_slide()
rect(s, 0, 0, EMU_W, EMU_H, NAVY)
rect(s, 0, Inches(4.55), EMU_W, Inches(0.07), TEAL)
tb, tf = textbox(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(2.0))
p = tf.paragraphs[0]
set_run(p.add_run(), "AI Bank Reconciliation Agent", 46, WHITE, bold=True)
p2 = tf.add_paragraph()
p2.space_before = Pt(14)
set_run(p2.add_run(), "Automated month-end reconciliation for your physician practices",
        22, RGBColor(0xBF, 0xD8, 0xD3))
tb, tf = textbox(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.5))
set_run(tf.paragraphs[0].add_run(), "Client overview", 14, TEAL, bold=True)
notes(s, "Today I'll walk you through an agent that takes over the monthly bank "
         "reconciliation for each practice — from collecting the documents all the "
         "way to the finished report you send the physician — while keeping your "
         "accountants in control of every decision that matters. I'll show you how "
         "it works, why you can trust the numbers, and how we'd roll it out.")

# ---- Slide 2: The problem today -------------------------------------------
s = add_slide()
header(s, "The problem today")
bullets(s, [
    "Every practice generates a stack of financial documents each month, from many sources and in every format",
    ("Bank statements, insurance remittances (ERAs/EOBs), EFT notices, government payor remittances", 1),
    ("PMS reports, GL exports, payroll, vendor invoices, check registers, utility bills", 1),
    "They arrive by email, shared drives, and manual upload",
    "An accountant manually reads, matches, chases discrepancies, and writes it up",
    "Slow, repetitive, and error-prone — and it repeats for every practice, every month",
])
notes(s, "This is work your team does today by hand. It's not hard in concept — match "
         "each deposit to the payment that explains it, match each withdrawal to its "
         "expense, resolve what doesn't line up, prove the bank agrees with the books. "
         "But it's hours of careful, repetitive effort, multiplied across every "
         "practice and every month.")

# ---- Slide 3: The big idea ------------------------------------------------
s = add_slide()
rect(s, 0, 0, EMU_W, EMU_H, LIGHT)
rect(s, 0, 0, Inches(0.25), EMU_H, TEAL)
tb, tf = textbox(s, Inches(0.9), Inches(0.6), Inches(11.5), Inches(0.5))
set_run(tf.paragraphs[0].add_run(), "THE ONE IDEA THAT MAKES THIS TRUSTWORTHY", 14, TEAL, bold=True)
panel = rect(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(2.5), NAVY)
tf = panel.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.5); tf.margin_right = Inches(0.5)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
set_run(p.add_run(), "The AI reads and understands the documents.", 28, WHITE, bold=True)
p2 = tf.add_paragraph()
set_run(p2.add_run(), "The math is done by audited, deterministic code — not the AI's guess.",
        28, RGBColor(0x7F, 0xD4, 0xC8), bold=True)
bullets(s, [
    "AI does what it's good at: reading messy paperwork, extracting numbers, explaining discrepancies, answering questions",
    "The matching and balancing is done by precise, repeatable rules",
    "Every number traces back to its exact source document — file, page, line",
], top=Inches(4.4), size=18)
notes(s, "This is the most important slide. A natural worry with AI in finance is, "
         "'Can I trust the numbers?' Our answer is architectural: the AI never "
         "decides that your books balance. It reads the documents and does the "
         "investigative work, but the arithmetic is done by deterministic code that "
         "produces the same answer every time. And every figure links back to the "
         "document it came from. So the report isn't a black box — it's fully "
         "explainable and auditable.")

# ---- Slide 4: How a month flows -------------------------------------------
s = add_slide()
header(s, "How a month flows through the agent")
steps = [
    ("1", "Documents arrive", "Email, drives, uploads — collected automatically"),
    ("2", "Read & classify", "Each document recognized; data extracted"),
    ("✋", "Confirm set complete", "You confirm the month is complete"),
    ("4", "Reconcile", "Deposits and withdrawals matched; balance proven"),
    ("5", "Investigate exceptions", "Conversationally — just ask"),
    ("✋", "Approve resolutions", "You approve adjustments"),
    ("✋", "Final sign-off", "Professional judgment stays with you"),
    ("8", "Excel report", "Generated, ready for the physician"),
]
x0, y0 = Inches(0.55), Inches(2.1)
cw, ch, gapx, gapy = Inches(3.0), Inches(1.7), Inches(0.18), Inches(0.35)
for i, (num, title, sub) in enumerate(steps):
    col = i % 4
    row = i // 4
    x = x0 + col * (cw + gapx)
    y = y0 + row * (ch + gapy)
    is_gate = num == "✋"
    card = rect(s, x, y, cw, ch, ACCENT_BG if is_gate else WHITE)
    card.line.color.rgb = TEAL if is_gate else RGBColor(0xD5, 0xDD, 0xE2)
    card.line.width = Pt(1.5 if is_gate else 0.75)
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.15); tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    set_run(p.add_run(), num + "   ", 16, TEAL if is_gate else NAVY, bold=True)
    set_run(p.add_run(), title, 15, NAVY, bold=True)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(4)
    set_run(p2.add_run(), sub, 11, MUTE)
notes(s, "Here's the whole journey on one slide. Notice the three hand-raise symbols — "
         "those are the points where the agent stops and waits for your team. "
         "Everything between them is automated. I'll come back to those checkpoints "
         "in a moment.")

# ---- Slide 5: Messy real world --------------------------------------------
s = add_slide()
header(s, "It handles the messy real world")
bullets(s, [
    "One insurance deposit that covers dozens of patient claims",
    "A card-processor batch that bundles many copays into one deposit",
    "Payor takebacks that reduce a deposit",
    "The few-day lag between a recorded deposit and when it clears the bank",
    "Bank fees, payroll, outstanding checks, and missing remittances",
], size=19)
notes(s, "Reconciliation isn't just matching equal amounts. A single Aetna deposit "
         "might cover forty claims. A card processor lumps a day's copays into one "
         "bank line. A payor claws back a prior overpayment, so the deposit is short "
         "by exactly that amount. The agent is built for these realities — it doesn't "
         "just give up when the numbers aren't a clean one-to-one match.")

# ---- Slide 6: Conversational example --------------------------------------
s = add_slide()
header(s, "You investigate by simply asking")
# user bubble
u = rect(s, Inches(0.7), Inches(2.0), Inches(8.5), Inches(0.95), NAVY)
tf = u.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.25); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
set_run(tf.paragraphs[0].add_run(), "You:  “Why is this $4,812 deposit unmatched?”",
        18, WHITE, italic=True)
# agent bubble
a = rect(s, Inches(2.5), Inches(3.15), Inches(10.1), Inches(1.5), ACCENT_BG)
a.line.color.rgb = TEAL; a.line.width = Pt(1.25)
tf = a.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.25); tf.margin_right = Inches(0.25); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
set_run(p.add_run(), "Agent:  ", 18, TEAL, bold=True)
set_run(p.add_run(), "“It's the Aetna ERA for $4,900, reduced by an $87.67 recoupment "
        "on claim #123. The remittance arrived on the 14th. Shall I record the match?”",
        18, INK, italic=True)
bullets(s, [
    "Plain-language investigation of any exception",
    "The agent proposes a resolution and explains its reasoning",
    "You accept, reject, or ask follow-ups",
], top=Inches(5.0), size=17)
notes(s, "When something doesn't tie out, you don't dig through spreadsheets — you ask. "
         "The agent investigates and comes back with an explanation and a proposed "
         "fix. You stay in the driver's seat; it does the legwork.")

# ---- Slide 7: Checkpoints table -------------------------------------------
s = add_slide()
header(s, "Where humans stay in control")
rows = [
    ("✋  Checkpoint", "Why it matters"),
    ("Confirm the month's documents are complete", "Don't reconcile on a partial picture"),
    ("Approve large or low-confidence matches", "A person verifies anything material or uncertain"),
    ("Approve exception resolutions", "Adjustments and write-offs are human decisions"),
    ("Final sign-off before the report goes out", "Professional judgment stays with the accountant"),
]
tbl_shape = s.shapes.add_table(len(rows), 2, Inches(0.7), Inches(2.0),
                               Inches(12.0), Inches(4.0))
table = tbl_shape.table
table.columns[0].width = Inches(6.2)
table.columns[1].width = Inches(5.8)
for r, (c0, c1) in enumerate(rows):
    for c, text in enumerate((c0, c1)):
        cell = table.cell(r, c)
        cell.margin_left = Inches(0.15); cell.margin_top = Inches(0.08); cell.margin_bottom = Inches(0.08)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        if r == 0:
            cell.fill.fore_color.rgb = NAVY
        else:
            cell.fill.fore_color.rgb = WHITE if r % 2 else ACCENT_BG
        para = cell.text_frame.paragraphs[0]
        set_run(para.add_run(), text, 15 if r == 0 else 14,
                WHITE if r == 0 else INK, bold=(r == 0 or c == 0))
notes(s, "The agent is deliberately designed to stop at the moments that carry "
         "professional or financial responsibility. It proposes; your team decides. "
         "This isn't 'AI replaces the accountant' — it's 'AI does the grunt work, the "
         "accountant keeps the judgment.'")

# ---- Slide 8: Trust -------------------------------------------------------
s = add_slide()
header(s, "Why you can trust the numbers")
cards = [
    ("Deterministic & repeatable", "Same inputs, same result, every time"),
    ("Fully traceable", "Every figure links to its source document"),
    ("Completely logged", "Every read, match, approval, and override is recorded"),
    ("Nothing dropped silently", "Unmatched items become visible exceptions"),
]
x0, y0 = Inches(0.7), Inches(2.2)
cw, ch, gx, gy = Inches(5.85), Inches(1.9), Inches(0.3), Inches(0.4)
for i, (t, d) in enumerate(cards):
    x = x0 + (i % 2) * (cw + gx)
    y = y0 + (i // 2) * (ch + gy)
    card = rect(s, x, y, cw, ch, WHITE)
    card.line.color.rgb = TEAL; card.line.width = Pt(1.0)
    rect(s, x, y, Inches(0.12), ch, TEAL)
    tf = card.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.2)
    set_run(tf.paragraphs[0].add_run(), t, 19, NAVY, bold=True)
    p = tf.add_paragraph(); p.space_before = Pt(6)
    set_run(p.add_run(), d, 14, MUTE)
notes(s, "Four guarantees. The math is repeatable, every number is traceable, every "
         "action is logged, and nothing falls through the cracks. That's what makes "
         "the report defensible in an audit and trustworthy to the physician who "
         "receives it.")

# ---- Slide 9: Healthcare data ---------------------------------------------
s = add_slide()
header(s, "Built for healthcare data")
bullets(s, [
    "Insurance remittances contain protected health information (PHI)",
    "Each practice fully isolated as its own project",
    "Documents encrypted; complete audit trail for compliance",
    "AI used under a HIPAA-compliant arrangement — hosting kept flexible to fit your infrastructure",
], size=19)
notes(s, "Because ERAs and EOBs contain PHI, this is built for healthcare from the "
         "ground up — isolation per practice, encryption, and a full audit trail. The "
         "exact AI hosting model is something we'll finalize with your compliance team "
         "so it fits what you already have.")

# ---- Slide 10: Deliverable ------------------------------------------------
s = add_slide()
header(s, "The deliverable")
tb, tf = textbox(s, Inches(0.7), Inches(1.85), Inches(12), Inches(0.5))
set_run(tf.paragraphs[0].add_run(),
        "A structured Excel workbook — the same report you send physicians today:",
        18, INK, bold=True)
tabs = ["Summary & sign-off", "Formal bank reconciliation", "Matched deposits",
        "Matched disbursements", "Outstanding checks", "Deposits in transit",
        "Exceptions & resolutions", "Source-document index"]
x0, y0 = Inches(0.7), Inches(2.7)
cw, ch, gx, gy = Inches(3.85), Inches(0.95), Inches(0.25), Inches(0.3)
for i, t in enumerate(tabs):
    x = x0 + (i % 3) * (cw + gx)
    y = y0 + (i // 3) * (ch + gy)
    card = rect(s, x, y, cw, ch, ACCENT_BG)
    card.line.color.rgb = TEAL; card.line.width = Pt(0.75)
    tf = card.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.2)
    set_run(tf.paragraphs[0].add_run(), t, 15, NAVY, bold=True)
notes(s, "The output is exactly what your accountant produces today — a clean Excel "
         "report — but with the matching done and every figure backed by its source. "
         "Your accountant reviews, signs off, and sends it on.")

# ---- Slide 11: Rollout ----------------------------------------------------
s = add_slide()
header(s, "How we'd roll it out")
phases = [
    ("1", "Pilot one practice, one month", "Run alongside the current manual process and confirm the output matches"),
    ("2", "Tune", "Adjust matching and checkpoints to that practice's payors and habits"),
    ("3", "Expand practice by practice", "Each isolated, reusing what the agent learns about payors and formats"),
]
y = Inches(2.1)
for num, t, d in phases:
    circ = rect(s, Inches(0.8), y, Inches(0.9), Inches(0.9), TEAL)
    circ.adjustments  # keep as rect-circle proxy
    tf = circ.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    pp = tf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
    set_run(pp.add_run(), num, 26, WHITE, bold=True)
    tb, tf = textbox(s, Inches(2.0), y, Inches(10.5), Inches(1.2))
    p = tf.paragraphs[0]
    set_run(p.add_run(), t, 20, NAVY, bold=True)
    p2 = tf.add_paragraph(); p2.space_before = Pt(3)
    set_run(p2.add_run(), d, 15, MUTE)
    y += Inches(1.5)
notes(s, "We'd start safe: one practice, one month, in parallel with your existing "
         "process, so you can verify the agent against a known-good result before "
         "relying on it. Then we tune and expand.")

# ---- Slide 12: One sentence -----------------------------------------------
s = add_slide()
rect(s, 0, 0, EMU_W, EMU_H, NAVY)
rect(s, Inches(0.9), Inches(2.2), Inches(0.15), Inches(3.0), TEAL)
tb, tf = textbox(s, Inches(1.3), Inches(2.0), Inches(11.0), Inches(3.5))
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
set_run(tf.paragraphs[0].add_run(), "In one sentence", 16, TEAL, bold=True)
p = tf.add_paragraph(); p.space_before = Pt(14)
set_run(p.add_run(),
        "A conversational agent that ingests every financial document a practice "
        "receives, reconciles the bank to the books with auditable precision, helps "
        "resolve every exception, pauses for your team's approval at the moments that "
        "matter, and delivers a finished month-end report ready for the physician.",
        24, WHITE, bold=True)
notes(s, "To sum up: it does the heavy lifting end-to-end, you keep control of the "
         "decisions, and you get back hours per practice every month. Happy to take "
         "questions or walk through a live example.")

# ---- Save -----------------------------------------------------------------
out = Path(__file__).resolve().parent.parent / "docs" / "reconciliation-agent-deck.pptx"
prs.save(str(out))
print(f"Wrote {out} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
